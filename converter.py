from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pathlib import Path
import tempfile

def extract_pptx_to_md(file_bytes: bytes, filename: str) -> str:
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp.write(file_bytes)
        tmp_path = Path(tmp.name)

    prs = Presentation(str(tmp_path))
    file_name = Path(filename).stem
    lines = [f"# {file_name}\n"]

    for i, slide in enumerate(prs.slides, start=1):
        shapes = list(slide.shapes)

        # Title
        for shape in shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(f"## {shape.text_frame.text.strip()} (Slide {i})\n")
                break

        # Body
        for shape in shapes[1:]:
            if shape.has_text_frame and shape.text_frame.text.strip():
                lines.append(f"{shape.text_frame.text.strip()}\n")

        # Tables
        for shape in shapes:
            if shape.has_table:
                table = shape.table
                for j, row in enumerate(table.rows):
                    cells = [cell.text.replace("\n", " ") for cell in row.cells]
                    lines.append("| " + " | ".join(cells) + " |")
                    if j == 0:
                        lines.append("| " + " | ".join(["---"] * len(cells)) + " |")
                lines.append("")

    tmp_path.unlink()
    return "\n".join(lines)